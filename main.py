import os
import re
from tkinter import Tk, Label, Button, filedialog, StringVar, IntVar, ttk
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def adjust_image_to_16_9(image_path):
    """調整圖片為 16:9 的比例"""
    img = Image.open(image_path)
    width, height = img.size

    # 計算目標比例
    target_ratio = 16 / 9
    current_ratio = width / height

    if current_ratio > target_ratio:  # 寬太長
        new_width = int(height * target_ratio)
        offset = (width - new_width) // 2
        cropped_img = img.crop((offset, 0, offset + new_width, height))
    else:  # 高太長
        new_height = int(width / target_ratio)
        offset = (height - new_height) // 2
        cropped_img = img.crop((0, offset, width, offset + new_height))

    return cropped_img


def sort_files_by_number(files):
    """根據檔案名稱中的數字大小排序"""
    def extract_number(filename):
        match = re.search(r'\d+', filename)
        return int(match.group()) if match else float('inf')  # 若無數字，置於最後
    return sorted(files, key=extract_number)


def add_images_to_ppt(folder_path, output_ppt, progress_var, progress_bar, progress_label):
    """將資料夾中的 PNG 圖片新增到 PPT 並設定為 16:9"""
    prs = Presentation()
    
    # 設定投影片比例為 16:9
    prs.slide_width = Inches(13.33)  # 寬度
    prs.slide_height = Inches(7.5)   # 高度 (16:9)

    # 獲取資料夾中的所有 PNG 圖片，並按檔案名稱中的數字大小排序
    png_files = [f for f in os.listdir(folder_path) if f.lower().endswith('.png')]
    png_files = sort_files_by_number(png_files)

    total_files = len(png_files)
    progress_var.set(0)
    progress_bar['maximum'] = total_files

    for index, file_name in enumerate(png_files):
        image_path = os.path.join(folder_path, file_name)
        
        # 調整圖片比例
        adjusted_img = adjust_image_to_16_9(image_path)
        adjusted_img_path = os.path.join(folder_path, f"adjusted_{file_name}")
        adjusted_img.save(adjusted_img_path)

        # 新增幻燈片
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白幻燈片

        # 設定圖片大小與位置
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height

        slide.shapes.add_picture(adjusted_img_path, left, top, width, height)
        os.remove(adjusted_img_path)  # 移除調整過的臨時圖片

        # 更新進度條與進度標籤
        progress_var.set(index + 1)
        progress_label.config(text=f"進度：{index + 1}/{total_files}")
        progress_bar.update()

    # 儲存簡報
    prs.save(output_ppt)
    progress_label.config(text="處理完成！")
    print(f"PPT 已儲存為 {output_ppt}")


def browse_input_folder():
    folder_selected = filedialog.askdirectory()
    input_folder.set(folder_selected)
    output_file.set(os.path.join(folder_selected, "output.pptx"))


def browse_output_file():
    file_selected = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                 filetypes=[("PowerPoint Files", "*.pptx")])
    output_file.set(file_selected)


def generate_ppt():
    folder_path = input_folder.get()
    output_ppt = output_file.get()
    if not folder_path:
        print("請選擇輸入資料夾！")
        return
    if not output_ppt:
        print("請選擇輸出檔案位置！")
        return

    # 顯示進度條
    progress_bar.grid(row=4, column=0, columnspan=3, pady=5, sticky="ew")
    progress_label.grid(row=5, column=0, columnspan=3, pady=5)
    add_images_to_ppt(folder_path, output_ppt, progress_var, progress_bar, progress_label)


# 建立 tkinter 介面
root = Tk()
root.title("圖片匯入 PPT 工具")

# 輸入資料夾
input_folder = StringVar()
output_file = StringVar()

Label(root, text="選擇輸入資料夾：").grid(row=0, column=0, padx=5, pady=5)
Button(root, text="瀏覽", command=browse_input_folder).grid(row=0, column=2, padx=5, pady=5)
Label(root, textvariable=input_folder, width=50).grid(row=0, column=1, padx=5, pady=5)

# 輸出檔案
Label(root, text="選擇輸出檔案位置：").grid(row=1, column=0, padx=5, pady=5)
Button(root, text="瀏覽", command=browse_output_file).grid(row=1, column=2, padx=5, pady=5)
Label(root, textvariable=output_file, width=50).grid(row=1, column=1, padx=5, pady=5)

# 開始處理按鈕
Button(root, text="開始處理", command=generate_ppt).grid(row=3, column=0, columnspan=3, pady=10)

# 隱藏的進度條與進度標籤
progress_var = IntVar()
progress_bar = ttk.Progressbar(root, variable=progress_var, maximum=100)
progress_label = Label(root, text="", font=("Arial", 10))

root.mainloop()
